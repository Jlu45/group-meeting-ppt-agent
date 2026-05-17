import os
import uuid
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu

from src.common.models import SlideSpec, TemplateDNA, ValidationIssue


class VisualValidator:
    MARGIN_TOLERANCE_EMU = 91440
    MIN_FONT_SIZE = Pt(8)
    MAX_FONT_SIZE = Pt(44)
    PLACEHOLDER_TEXTS = {
        "单击此处添加标题", "单击此处添加文本", "Click to add title",
        "Click to add text", "标题", "文本", "Title", "Text",
    }

    def validate(self, pptx_path: str, template_dna: TemplateDNA = None,
                 slide_specs: list[SlideSpec] = None) -> list[ValidationIssue]:
        prs = Presentation(pptx_path)
        issues = []
        issues.extend(self._validate_structure(prs))
        issues.extend(self._validate_layout(prs))
        if template_dna:
            issues.extend(self._validate_compliance(prs, template_dna))
        if slide_specs:
            issues.extend(self._validate_content(prs, slide_specs))
        return issues

    def _validate_structure(self, prs) -> list[ValidationIssue]:
        issues = []
        slide_count = len(prs.slides)
        if slide_count == 0:
            issues.append(ValidationIssue(
                id=str(uuid.uuid4())[:8],
                severity="critical",
                issue_type="structure",
                message="演示文稿不包含任何幻灯片",
                suggested_fix="添加至少一张幻灯片",
            ))
            return issues

        if slide_count > 50:
            issues.append(ValidationIssue(
                id=str(uuid.uuid4())[:8],
                severity="warning",
                issue_type="structure",
                message=f"幻灯片数量过多({slide_count}张)，建议控制在30张以内",
                suggested_fix="精简内容，合并相关幻灯片",
            ))

        for idx, slide in enumerate(prs.slides):
            slide_id = f"slide_{idx + 1}"
            has_text = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            has_text = True
                            break
                if has_text:
                    break
            if not has_text and len(slide.shapes) <= 1:
                issues.append(ValidationIssue(
                    id=str(uuid.uuid4())[:8],
                    severity="warning",
                    slide_id=slide_id,
                    issue_type="blank_slide",
                    message=f"第{idx + 1}张幻灯片为空白页",
                    suggested_fix="添加内容或删除该幻灯片",
                ))

            has_title = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape.shape_type == 14:
                        has_title = True
                        break
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                        if shape.placeholder_format.type == 1:
                            has_title = True
                            break
            if not has_title and idx > 0:
                issues.append(ValidationIssue(
                    id=str(uuid.uuid4())[:8],
                    severity="warning",
                    slide_id=slide_id,
                    issue_type="missing_title",
                    message=f"第{idx + 1}张幻灯片缺少标题",
                    suggested_fix="添加标题占位符",
                ))

        return issues

    def _validate_layout(self, prs) -> list[ValidationIssue]:
        issues = []
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for idx, slide in enumerate(prs.slides):
            slide_id = f"slide_{idx + 1}"
            shape_rects = []

            for shape in slide.shapes:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                right = left + width
                bottom = top + height

                if left < -self.MARGIN_TOLERANCE_EMU:
                    issues.append(ValidationIssue(
                        id=str(uuid.uuid4())[:8],
                        severity="warning",
                        slide_id=slide_id,
                        element_id=shape.shape_id,
                        issue_type="margin_violation",
                        message=f"第{idx + 1}张幻灯片元素超出左边界",
                        suggested_fix="调整元素位置使其在幻灯片边界内",
                    ))
                if top < -self.MARGIN_TOLERANCE_EMU:
                    issues.append(ValidationIssue(
                        id=str(uuid.uuid4())[:8],
                        severity="warning",
                        slide_id=slide_id,
                        element_id=shape.shape_id,
                        issue_type="margin_violation",
                        message=f"第{idx + 1}张幻灯片元素超出上边界",
                        suggested_fix="调整元素位置使其在幻灯片边界内",
                    ))
                if right > slide_width + self.MARGIN_TOLERANCE_EMU:
                    issues.append(ValidationIssue(
                        id=str(uuid.uuid4())[:8],
                        severity="warning",
                        slide_id=slide_id,
                        element_id=shape.shape_id,
                        issue_type="margin_violation",
                        message=f"第{idx + 1}张幻灯片元素超出右边界",
                        suggested_fix="调整元素宽度或位置",
                    ))
                if bottom > slide_height + self.MARGIN_TOLERANCE_EMU:
                    issues.append(ValidationIssue(
                        id=str(uuid.uuid4())[:8],
                        severity="warning",
                        slide_id=slide_id,
                        element_id=shape.shape_id,
                        issue_type="margin_violation",
                        message=f"第{idx + 1}张幻灯片元素超出下边界",
                        suggested_fix="调整元素高度或位置",
                    ))

                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text_len = len(para.text)
                        if text_len > 200:
                            issues.append(ValidationIssue(
                                id=str(uuid.uuid4())[:8],
                                severity="warning",
                                slide_id=slide_id,
                                element_id=shape.shape_id,
                                issue_type="text_overflow",
                                message=f"第{idx + 1}张幻灯片文本段落过长({text_len}字符)",
                                suggested_fix="缩短文本或拆分为多个要点",
                            ))
                        for run in para.runs:
                            if run.font.size and run.font.size < self.MIN_FONT_SIZE:
                                issues.append(ValidationIssue(
                                    id=str(uuid.uuid4())[:8],
                                    severity="warning",
                                    slide_id=slide_id,
                                    element_id=shape.shape_id,
                                    issue_type="font_too_small",
                                    message=f"第{idx + 1}张幻灯片字体过小({run.font.size})",
                                    suggested_fix="将字体大小调整为至少8pt",
                                ))

                shape_rects.append((left, top, right, bottom, shape.shape_id))

            for i in range(len(shape_rects)):
                for j in range(i + 1, len(shape_rects)):
                    r1 = shape_rects[i]
                    r2 = shape_rects[j]
                    overlap_x = max(0, min(r1[2], r2[2]) - max(r1[0], r2[0]))
                    overlap_y = max(0, min(r1[3], r2[3]) - max(r1[1], r2[1]))
                    area1 = (r1[2] - r1[0]) * (r1[3] - r1[1])
                    area2 = (r2[2] - r2[0]) * (r2[3] - r2[1])
                    overlap_area = overlap_x * overlap_y
                    min_area = min(area1, area2)
                    if min_area > 0 and overlap_area / min_area > 0.3:
                        issues.append(ValidationIssue(
                            id=str(uuid.uuid4())[:8],
                            severity="warning",
                            slide_id=slide_id,
                            element_id=f"{r1[4]}_{r2[4]}",
                            issue_type="element_overlap",
                            message=f"第{idx + 1}张幻灯片元素重叠超过30%",
                            suggested_fix="调整元素位置避免重叠",
                        ))

        return issues

    def _validate_compliance(self, prs, template_dna) -> list[ValidationIssue]:
        issues = []
        theme_colors = {}
        theme_fonts = {}
        if template_dna and template_dna.theme:
            theme_colors = template_dna.theme.colors or {}
            theme_fonts = template_dna.theme.fonts or {}

        allowed_colors = set()
        for color_val in theme_colors.values():
            if isinstance(color_val, str):
                allowed_colors.add(color_val.upper().replace("#", ""))

        allowed_font_names = set()
        for font_val in theme_fonts.values():
            if isinstance(font_val, str):
                allowed_font_names.add(font_val)

        for idx, slide in enumerate(prs.slides):
            slide_id = f"slide_{idx + 1}"
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name and allowed_font_names and run.font.name not in allowed_font_names:
                            issues.append(ValidationIssue(
                                id=str(uuid.uuid4())[:8],
                                severity="info",
                                slide_id=slide_id,
                                element_id=shape.shape_id,
                                issue_type="font_mismatch",
                                message=f"第{idx + 1}张幻灯片字体'{run.font.name}'与模板不匹配",
                                suggested_fix=f"使用模板字体: {', '.join(allowed_font_names)}",
                            ))
                        try:
                            if run.font.color and run.font.color.type is not None:
                                color_str = str(run.font.color.rgb).upper()
                                if allowed_colors and color_str not in allowed_colors:
                                    issues.append(ValidationIssue(
                                        id=str(uuid.uuid4())[:8],
                                        severity="info",
                                        slide_id=slide_id,
                                        element_id=shape.shape_id,
                                        issue_type="color_mismatch",
                                        message=f"第{idx + 1}张幻灯片颜色#{color_str}与模板不匹配",
                                        suggested_fix=f"使用模板颜色: {', '.join(allowed_colors)}",
                                    ))
                        except (AttributeError, TypeError):
                            pass

        return issues

    def _validate_content(self, prs, slide_specs) -> list[ValidationIssue]:
        issues = []
        spec_map = {}
        if slide_specs:
            for spec in slide_specs:
                if spec.id:
                    spec_map[spec.id] = spec

        for idx, slide in enumerate(prs.slides):
            slide_id = f"slide_{idx + 1}"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text in self.PLACEHOLDER_TEXTS:
                            issues.append(ValidationIssue(
                                id=str(uuid.uuid4())[:8],
                                severity="warning",
                                slide_id=slide_id,
                                element_id=shape.shape_id,
                                issue_type="placeholder_text",
                                message=f"第{idx + 1}张幻灯片包含占位符文本'{text}'",
                                suggested_fix="替换为实际内容",
                            ))

        if slide_specs and len(prs.slides) < len(slide_specs):
            issues.append(ValidationIssue(
                id=str(uuid.uuid4())[:8],
                severity="warning",
                issue_type="incomplete_content",
                message=f"生成的幻灯片({len(prs.slides)}张)少于规划({len(slide_specs)}张)",
                suggested_fix="检查是否有内容在生成过程中丢失",
            ))

        return issues

    def auto_fix(self, pptx_path: str, issues: list[ValidationIssue],
                 max_rounds: int = 3) -> tuple[str, list[ValidationIssue]]:
        base, ext = os.path.splitext(pptx_path)
        fixed_path = f"{base}_fixed{ext}"
        import shutil
        shutil.copy2(pptx_path, fixed_path)

        remaining = list(issues)
        for round_num in range(max_rounds):
            if not remaining:
                break

            prs = Presentation(fixed_path)
            fixed_in_round = set()

            if round_num == 0:
                for issue in remaining:
                    if issue.issue_type == "text_overflow":
                        self._fix_overflow(prs, issue)
                        fixed_in_round.add(issue.id)
                    elif issue.issue_type == "blank_slide":
                        self._fix_blank_slide(prs, issue)
                        fixed_in_round.add(issue.id)
            elif round_num == 1:
                for issue in remaining:
                    if issue.issue_type in ("font_mismatch", "font_too_small"):
                        self._fix_font(prs, issue)
                        fixed_in_round.add(issue.id)
                    elif issue.issue_type == "color_mismatch":
                        self._fix_color(prs, issue)
                        fixed_in_round.add(issue.id)
            else:
                for issue in remaining:
                    if issue.issue_type == "margin_violation":
                        self._fix_margin(prs, issue)
                        fixed_in_round.add(issue.id)
                    elif issue.issue_type == "placeholder_text":
                        self._fix_placeholder(prs, issue)
                        fixed_in_round.add(issue.id)

            prs.save(fixed_path)
            remaining = [i for i in remaining if i.id not in fixed_in_round]

        return fixed_path, remaining

    def _fix_overflow(self, prs, issue):
        try:
            slide_idx = int(issue.slide_id.split("_")[1]) - 1
            slide = prs.slides[slide_idx]
            for shape in slide.shapes:
                if shape.has_text_frame and str(shape.shape_id) == str(issue.element_id):
                    for para in shape.text_frame.paragraphs:
                        if len(para.text) > 200:
                            for run in para.runs:
                                if len(run.text) > 80:
                                    run.text = run.text[:79] + "…"
        except (IndexError, ValueError, AttributeError):
            pass

    def _fix_blank_slide(self, prs, issue):
        try:
            slide_idx = int(issue.slide_id.split("_")[1]) - 1
            slide = prs.slides[slide_idx]
            if len(slide.shapes) == 0:
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                height = Inches(2)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = "（待补充内容）"
        except (IndexError, ValueError):
            pass

    def _fix_font(self, prs, issue):
        try:
            slide_idx = int(issue.slide_id.split("_")[1]) - 1
            slide = prs.slides[slide_idx]
            for shape in slide.shapes:
                if shape.has_text_frame and str(shape.shape_id) == str(issue.element_id):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if issue.issue_type == "font_too_small":
                                run.font.size = Pt(10)
                            elif issue.issue_type == "font_mismatch":
                                run.font.name = "微软雅黑"
        except (IndexError, ValueError, AttributeError):
            pass

    def _fix_color(self, prs, issue):
        try:
            slide_idx = int(issue.slide_id.split("_")[1]) - 1
            slide = prs.slides[slide_idx]
            for shape in slide.shapes:
                if shape.has_text_frame and str(shape.shape_id) == str(issue.element_id):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.color and run.font.color.rgb:
                                from pptx.dml.color import RGBColor
                                run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
        except (IndexError, ValueError, AttributeError):
            pass

    def _fix_margin(self, prs, issue):
        try:
            slide_idx = int(issue.slide_id.split("_")[1]) - 1
            slide = prs.slides[slide_idx]
            for shape in slide.shapes:
                if str(shape.shape_id) == str(issue.element_id):
                    if shape.left < 0:
                        shape.left = 0
                    if shape.top < 0:
                        shape.top = 0
                    if shape.left + shape.width > prs.slide_width:
                        shape.left = prs.slide_width - shape.width
                    if shape.top + shape.height > prs.slide_height:
                        shape.top = prs.slide_height - shape.height
        except (IndexError, ValueError, AttributeError):
            pass

    def _fix_placeholder(self, prs, issue):
        try:
            slide_idx = int(issue.slide_id.split("_")[1]) - 1
            slide = prs.slides[slide_idx]
            for shape in slide.shapes:
                if shape.has_text_frame and str(shape.shape_id) == str(issue.element_id):
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip() in self.PLACEHOLDER_TEXTS:
                            for run in para.runs:
                                if run.text.strip() in self.PLACEHOLDER_TEXTS:
                                    run.text = ""
        except (IndexError, ValueError, AttributeError):
            pass
