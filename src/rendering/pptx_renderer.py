import logging
import os
import tempfile
from typing import Any, Dict, List, Optional

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from src.common.models import (
    AssetStore,
    ChartIntent,
    DecorationSpec,
    LayoutSpec,
    PlaceholderSpec,
    Rect,
    SlideElementSpec,
    SlideSpec,
    TemplateDNA,
    ThemeSpec,
)
from src.rendering.layout_matcher import LayoutMatcher
from src.rendering.placeholder_binder import PlaceholderBinder
from src.template.layout_classifier import classify_layout

logger = logging.getLogger(__name__)

_DEFAULT_COLORS = {
    "primary": "#2B579A",
    "secondary": "#5B9BD5",
    "accent": "#ED7D31",
    "text": "#333333",
    "background": "#FFFFFF",
    "title": "#2B579A",
}

_DEFAULT_FONTS = {
    "title": "Calibri",
    "body": "Calibri",
}

_LAYOUT_TYPE_NAMES = {
    "cover": "Title Slide",
    "section": "Section Header",
    "bullet_list": "Title and Content",
    "two_column": "Two Content",
    "image_focus": "Picture with Caption",
    "chart": "Title and Content",
    "table": "Title and Content",
    "blank": "Blank",
    "generic": "Title and Content",
}

_PH_TYPE_ENUM_MAP = {
    "title": PP_PLACEHOLDER.TITLE,
    "ctrTitle": PP_PLACEHOLDER.CENTER_TITLE,
    "subtitle": PP_PLACEHOLDER.SUBTITLE,
    "body": PP_PLACEHOLDER.BODY,
    "obj": PP_PLACEHOLDER.OBJECT,
    "chart": PP_PLACEHOLDER.CHART,
    "pic": PP_PLACEHOLDER.PICTURE,
    "tbl": PP_PLACEHOLDER.TABLE,
    "ftr": PP_PLACEHOLDER.FOOTER,
    "sldNum": PP_PLACEHOLDER.SLIDE_NUMBER,
}


class LayoutDrivenRenderer:
    def __init__(self, template_dna: TemplateDNA):
        self.template_dna = template_dna
        self.layout_matcher = LayoutMatcher()
        self.placeholder_binder = PlaceholderBinder()
        self._chart_generator = None
        self._asset_store: Optional[AssetStore] = None

    def render(self, slide_specs: list[SlideSpec], output_path: str, asset_store: AssetStore = None) -> dict:
        self._asset_store = asset_store or AssetStore()

        template_path = (
            self.template_dna.media.get("template_path", "")
            or self.template_dna.media.get("source_path", "")
        )

        if template_path and os.path.isfile(template_path):
            prs = Presentation(template_path)
            if len(prs.slide_layouts) > 0:
                self._remove_existing_slides(prs)
        else:
            prs = Presentation()

        render_log: Dict[str, Any] = {
            "slide_count": 0,
            "layout_usage": {},
            "warnings": [],
            "errors": [],
        }

        layouts = self.template_dna.layouts if self.template_dna.layouts else self._build_fallback_layouts()

        for slide_spec in slide_specs:
            try:
                layout_spec = self.layout_matcher.select(slide_spec.intent, layouts)
                binding = self.placeholder_binder.bind(slide_spec, layout_spec)
                slide = self._create_slide_from_layout(prs, layout_spec)
                self._fill_placeholder(slide, binding, slide_spec, layout_spec)
                self._apply_decorations(slide, layout_spec, self.template_dna)

                if slide_spec.speaker_notes:
                    try:
                        notes_slide = slide.notes_slide
                        notes_slide.notes_text_frame.text = slide_spec.speaker_notes
                    except Exception:
                        pass

                slide_spec.selected_layout_id = layout_spec.id
                render_log["slide_count"] += 1
                layout_id = layout_spec.id
                render_log["layout_usage"][layout_id] = render_log["layout_usage"].get(layout_id, 0) + 1
            except Exception as e:
                msg = f"Slide '{slide_spec.title}': {e}"
                logger.error(msg)
                render_log["errors"].append(msg)

        out_dir = os.path.dirname(os.path.abspath(output_path))
        os.makedirs(out_dir, exist_ok=True)
        prs.save(output_path)

        return render_log

    def _create_slide_from_layout(self, prs, layout_spec: LayoutSpec):
        target_layout = None

        for slide_layout in prs.slide_layouts:
            if slide_layout.name == layout_spec.name:
                target_layout = slide_layout
                break

        if target_layout is None:
            layout_type = layout_spec.layout_type or classify_layout(layout_spec)
            expected_name = _LAYOUT_TYPE_NAMES.get(layout_type, "")
            if expected_name:
                for slide_layout in prs.slide_layouts:
                    if expected_name.lower() in slide_layout.name.lower():
                        target_layout = slide_layout
                        break

        if target_layout is None:
            for slide_layout in prs.slide_layouts:
                if len(list(slide_layout.placeholders)) > 0:
                    target_layout = slide_layout
                    break

        if target_layout is None:
            target_layout = prs.slide_layouts[0]

        return prs.slides.add_slide(target_layout)

    def _fill_placeholder(self, slide, binding: dict, slide_spec: SlideSpec, layout_spec: LayoutSpec):
        ph_specs = {p.id: p for p in layout_spec.placeholders}

        for element_id, placeholder_id in binding.items():
            ph_spec = ph_specs.get(placeholder_id)
            if ph_spec is None:
                continue

            ph_shape = self._find_placeholder_shape(slide, ph_spec)
            content = self._resolve_content(element_id, slide_spec)
            if content is None:
                continue

            if ph_shape is None:
                if isinstance(content, str):
                    self._add_text_box(slide, ph_spec, content)
                elif isinstance(content, dict) and "points" in content:
                    points = content["points"]
                    text = "\n".join(str(p) for p in points) if isinstance(points, list) else str(points)
                    self._add_text_box(slide, ph_spec, text)
                continue

            if isinstance(content, dict):
                if "chart_intent" in content:
                    ci_data = content["chart_intent"]
                    chart_intent = ChartIntent.from_dict(ci_data) if isinstance(ci_data, dict) else ci_data
                    image_path = self._generate_chart_image(chart_intent, self._asset_store)
                    self._add_image_to_placeholder(slide, ph_shape, ph_spec, image_path)
                elif "image_path" in content:
                    self._add_image_to_placeholder(slide, ph_shape, ph_spec, content["image_path"])
                elif "points" in content:
                    points = content["points"]
                    text = "\n".join(str(p) for p in points) if isinstance(points, list) else str(points)
                    self._set_text_content(ph_shape, text, ph_spec)
                else:
                    self._set_text_content(ph_shape, str(content), ph_spec)
            elif isinstance(content, ChartIntent):
                image_path = self._generate_chart_image(content, self._asset_store)
                self._add_image_to_placeholder(slide, ph_shape, ph_spec, image_path)
            else:
                text = str(content)
                if ph_spec.max_chars and len(text) > ph_spec.max_chars:
                    text = text[:max(ph_spec.max_chars - 3, 0)] + "..."
                self._set_text_content(ph_shape, text, ph_spec)

    def _apply_text_format(self, run, text_style: dict, theme: ThemeSpec):
        font_name = text_style.get("font_name") or theme.fonts.get("body") or _DEFAULT_FONTS.get("body")
        if font_name:
            run.font.name = font_name

        font_size = text_style.get("font_size")
        if font_size:
            run.font.size = Pt(font_size)

        font_color = text_style.get("font_color") or theme.colors.get("text")
        if font_color:
            rgb = self._hex_to_rgb(font_color)
            if rgb:
                run.font.color.rgb = rgb

        if "bold" in text_style and text_style["bold"] is not None:
            run.font.bold = text_style["bold"]

        if "italic" in text_style and text_style["italic"] is not None:
            run.font.italic = text_style["italic"]

    def _generate_chart_image(self, chart_intent: ChartIntent, asset_store: AssetStore) -> str:
        chart_type = chart_intent.preferred_chart_type or chart_intent.intent_type or "bar"
        theme_colors = self._get_theme_colors()
        color_values = list(theme_colors.values())

        data = self._load_chart_data(chart_intent, asset_store)
        labels = self._load_chart_labels(chart_intent, asset_store)

        fig, ax = plt.subplots(figsize=(8, 5))
        bg = theme_colors.get("background", "#FFFFFF")
        fig.patch.set_facecolor(bg)
        ax.set_facecolor(bg)

        try:
            if chart_type in ("bar", "column"):
                x = range(len(data))
                ax.bar(x, data, color=color_values[:len(data)])
                if labels:
                    ax.set_xticks(list(x))
                    ax.set_xticklabels(labels, rotation=45, ha="right")
            elif chart_type == "line":
                x = range(len(data))
                ax.plot(
                    list(x), data,
                    color=theme_colors.get("primary", "#2B579A"),
                    marker="o", linewidth=2,
                )
                if labels:
                    ax.set_xticks(list(x))
                    ax.set_xticklabels(labels, rotation=45, ha="right")
            elif chart_type == "pie":
                ax.pie(
                    data,
                    labels=labels if labels else None,
                    colors=color_values[:len(data)],
                    autopct="%1.1f%%",
                )
            elif chart_type == "scatter":
                x = range(len(data))
                ax.scatter(list(x), data, color=theme_colors.get("primary", "#2B579A"), s=60)
                if labels:
                    ax.set_xticks(list(x))
                    ax.set_xticklabels(labels, rotation=45, ha="right")
            else:
                x = range(len(data))
                ax.bar(x, data, color=color_values[:len(data)])
                if labels:
                    ax.set_xticks(list(x))
                    ax.set_xticklabels(labels, rotation=45, ha="right")
        except Exception as e:
            logger.warning(f"Chart generation error: {e}")
            ax.bar(range(len(data)), data, color=color_values[:len(data)])

        if chart_intent.title:
            ax.set_title(chart_intent.title, color=theme_colors.get("title", "#333333"), fontsize=14)

        if chart_type != "pie":
            ax.tick_params(colors=theme_colors.get("text", "#333333"))
            for spine in ax.spines.values():
                spine.set_color(theme_colors.get("text", "#CCCCCC"))

        plt.tight_layout()

        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        fig.savefig(tmp.name, dpi=150, bbox_inches="tight", facecolor=fig.get_facecolor())
        plt.close(fig)
        tmp.close()

        return tmp.name

    def _apply_decorations(self, slide, layout_spec: LayoutSpec, template_dna: TemplateDNA):
        all_decorations = list(layout_spec.decorations) + list(template_dna.decorations)
        if not all_decorations:
            return

        for dec in all_decorations:
            try:
                if dec.shape_type == "line":
                    self._add_divider_line(slide, dec)
                elif dec.shape_type in ("text", "textBox"):
                    self._add_decoration_text(slide, dec, template_dna.theme)
                elif dec.shape_type == "picture" and dec.media_rel_id:
                    pass
            except Exception as e:
                logger.warning(f"Decoration apply failed for {dec.id}: {e}")

    def _remove_existing_slides(self, prs):
        try:
            sldIdLst = prs.slides._sldIdLst
            for sldId in list(sldIdLst):
                rId = sldId.get(qn("r:id"))
                if rId:
                    prs.part.drop_rel(rId)
                sldIdLst.remove(sldId)
        except Exception as e:
            logger.warning(f"Failed to remove existing slides: {e}")

    def _find_placeholder_shape(self, slide, ph_spec: PlaceholderSpec):
        if ph_spec.idx is not None:
            for shape in slide.placeholders:
                try:
                    if shape.placeholder_format.idx == ph_spec.idx:
                        return shape
                except Exception:
                    continue

        expected_type = _PH_TYPE_ENUM_MAP.get(ph_spec.ph_type)
        if expected_type is not None:
            for shape in slide.placeholders:
                try:
                    if shape.placeholder_format.type == expected_type:
                        return shape
                except Exception:
                    continue

        return None

    def _resolve_content(self, element_id: str, slide_spec: SlideSpec):
        if element_id == "title" and slide_spec.title:
            return slide_spec.title

        for elem in slide_spec.elements:
            if elem.role == element_id:
                if elem.content:
                    return elem.content
                if elem.asset_ids and self._asset_store:
                    return self._resolve_asset_content(elem)

        return None

    def _resolve_asset_content(self, elem: SlideElementSpec):
        for asset_id in elem.asset_ids:
            figure = self._asset_store.figures.get(asset_id)
            if figure is not None and figure.path and os.path.isfile(figure.path):
                return {"image_path": figure.path}

            table = self._asset_store.tables.get(asset_id)
            if table is not None:
                return {"table_id": asset_id, "title": table.title}

        return None

    def _set_text_content(self, ph_shape, text: str, ph_spec: PlaceholderSpec):
        if not ph_shape.has_text_frame:
            return

        tf = ph_shape.text_frame
        tf.clear()

        lines = text.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            run = p.add_run()
            run.text = line
            self._apply_text_format(run, ph_spec.text_style, self.template_dna.theme)

    def _add_text_box(self, slide, ph_spec: PlaceholderSpec, text: str):
        left = int(ph_spec.rect.x) if ph_spec.rect.x else Emu(Inches(0.5))
        top = int(ph_spec.rect.y) if ph_spec.rect.y else Emu(Inches(1))
        width = int(ph_spec.rect.w) if ph_spec.rect.w else Emu(Inches(9))
        height = int(ph_spec.rect.h) if ph_spec.rect.h else Emu(Inches(1))

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        lines = text.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            run = p.add_run()
            run.text = line
            self._apply_text_format(run, ph_spec.text_style, self.template_dna.theme)

    def _add_image_to_placeholder(self, slide, ph_shape, ph_spec: PlaceholderSpec, image_path: str):
        if not os.path.isfile(image_path):
            logger.warning(f"Image not found: {image_path}")
            return

        left = int(ph_spec.rect.x) if ph_spec.rect.x else ph_shape.left
        top = int(ph_spec.rect.y) if ph_spec.rect.y else ph_shape.top
        width = int(ph_spec.rect.w) if ph_spec.rect.w else ph_shape.width
        height = int(ph_spec.rect.h) if ph_spec.rect.h else ph_shape.height

        slide.shapes.add_picture(image_path, left, top, width, height)

    def _hex_to_rgb(self, hex_str: str) -> Optional[RGBColor]:
        try:
            hex_str = hex_str.lstrip("#")
            if len(hex_str) == 3:
                hex_str = hex_str[0] * 2 + hex_str[1] * 2 + hex_str[2] * 2
            if len(hex_str) != 6:
                return None
            return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
        except Exception:
            return None

    def _get_theme_colors(self) -> dict:
        colors = dict(_DEFAULT_COLORS)
        if self.template_dna.theme.colors:
            colors.update(self.template_dna.theme.colors)
        return colors

    def _load_chart_data(self, chart_intent: ChartIntent, asset_store: AssetStore) -> list:
        for asset_id in chart_intent.data_asset_ids:
            table = asset_store.tables.get(asset_id)
            if table is not None:
                try:
                    import pandas as pd
                    if table.dataframe_ref and os.path.isfile(table.dataframe_ref):
                        df = pd.read_csv(table.dataframe_ref)
                        if chart_intent.y_fields:
                            return df[chart_intent.y_fields[0]].values.tolist()
                        if len(df.columns) > 1:
                            return df.iloc[:, 1].values.tolist()
                        return [1] * len(df)
                except Exception:
                    pass

        return [1, 2, 3, 4, 5]

    def _load_chart_labels(self, chart_intent: ChartIntent, asset_store: AssetStore) -> list:
        for asset_id in chart_intent.data_asset_ids:
            table = asset_store.tables.get(asset_id)
            if table is not None:
                try:
                    import pandas as pd
                    if table.dataframe_ref and os.path.isfile(table.dataframe_ref):
                        df = pd.read_csv(table.dataframe_ref)
                        if chart_intent.x_field:
                            return df[chart_intent.x_field].values.tolist()
                        if len(df.columns) > 0:
                            return df.iloc[:, 0].values.tolist()
                        return []
                except Exception:
                    pass

        return []

    def _add_divider_line(self, slide, dec: DecorationSpec):
        left = int(dec.rect.x) if dec.rect.x else 0
        top = int(dec.rect.y) if dec.rect.y else 0
        width = int(dec.rect.w) if dec.rect.w else int(slide.width)
        height = int(dec.rect.h) if dec.rect.h else 20000

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        accent_color = self.template_dna.theme.colors.get("accent", "#ED7D31")
        rgb = self._hex_to_rgb(accent_color)
        if rgb:
            shape.fill.fore_color.rgb = rgb
        shape.line.fill.background()

    def _add_decoration_text(self, slide, dec: DecorationSpec, theme: ThemeSpec):
        left = int(dec.rect.x) if dec.rect.x else Emu(Inches(0.5))
        top = int(dec.rect.y) if dec.rect.y else int(slide.height) - Emu(Inches(0.4))
        width = int(dec.rect.w) if dec.rect.w else Emu(Inches(9))
        height = int(dec.rect.h) if dec.rect.h else Emu(Inches(0.3))

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = ""
        self._apply_text_format(
            run,
            {"font_size": 9, "font_color": theme.colors.get("text", "#999999")},
            theme,
        )

    def _build_fallback_layouts(self) -> list[LayoutSpec]:
        return [
            LayoutSpec(
                id="fallback_cover",
                name="Title Slide",
                layout_type="cover",
                placeholders=[
                    PlaceholderSpec(
                        id="fallback_cover_title",
                        layout_id="fallback_cover",
                        ph_type="ctrTitle",
                        idx=0,
                        rect=Rect(x=457200, y=1828800, w=8229600, h=1143000),
                        text_style={"font_name": "Calibri", "font_size": 36, "font_color": "#2B579A", "bold": True},
                        max_chars=80,
                    ),
                    PlaceholderSpec(
                        id="fallback_cover_subtitle",
                        layout_id="fallback_cover",
                        ph_type="subtitle",
                        idx=1,
                        rect=Rect(x=457200, y=3200400, w=8229600, h=685800),
                        text_style={"font_name": "Calibri", "font_size": 18, "font_color": "#666666"},
                        max_chars=120,
                    ),
                ],
                tags=["cover"],
            ),
            LayoutSpec(
                id="fallback_section",
                name="Section Header",
                layout_type="section",
                placeholders=[
                    PlaceholderSpec(
                        id="fallback_section_title",
                        layout_id="fallback_section",
                        ph_type="ctrTitle",
                        idx=0,
                        rect=Rect(x=457200, y=1371600, w=8229600, h=1143000),
                        text_style={"font_name": "Calibri", "font_size": 32, "font_color": "#2B579A", "bold": True},
                        max_chars=80,
                    ),
                    PlaceholderSpec(
                        id="fallback_section_subtitle",
                        layout_id="fallback_section",
                        ph_type="subtitle",
                        idx=1,
                        rect=Rect(x=457200, y=2743200, w=8229600, h=685800),
                        text_style={"font_name": "Calibri", "font_size": 18, "font_color": "#666666"},
                        max_chars=120,
                    ),
                ],
                tags=["section"],
            ),
            LayoutSpec(
                id="fallback_content",
                name="Title and Content",
                layout_type="bullet_list",
                placeholders=[
                    PlaceholderSpec(
                        id="fallback_content_title",
                        layout_id="fallback_content",
                        ph_type="title",
                        idx=0,
                        rect=Rect(x=457200, y=274638, w=8229600, h=1143000),
                        text_style={"font_name": "Calibri", "font_size": 28, "font_color": "#2B579A", "bold": True},
                        max_chars=60,
                    ),
                    PlaceholderSpec(
                        id="fallback_content_body",
                        layout_id="fallback_content",
                        ph_type="body",
                        idx=1,
                        rect=Rect(x=457200, y=1600200, w=8229600, h=4572000),
                        text_style={"font_name": "Calibri", "font_size": 16, "font_color": "#333333"},
                        max_chars=500,
                    ),
                ],
                tags=["content"],
            ),
            LayoutSpec(
                id="fallback_two_col",
                name="Two Content",
                layout_type="two_column",
                placeholders=[
                    PlaceholderSpec(
                        id="fallback_two_col_title",
                        layout_id="fallback_two_col",
                        ph_type="title",
                        idx=0,
                        rect=Rect(x=457200, y=274638, w=8229600, h=1143000),
                        text_style={"font_name": "Calibri", "font_size": 28, "font_color": "#2B579A", "bold": True},
                        max_chars=60,
                    ),
                    PlaceholderSpec(
                        id="fallback_two_col_body1",
                        layout_id="fallback_two_col",
                        ph_type="body",
                        idx=1,
                        rect=Rect(x=457200, y=1600200, w=3886200, h=4572000),
                        text_style={"font_name": "Calibri", "font_size": 14, "font_color": "#333333"},
                        max_chars=300,
                    ),
                    PlaceholderSpec(
                        id="fallback_two_col_body2",
                        layout_id="fallback_two_col",
                        ph_type="body",
                        idx=2,
                        rect=Rect(x=4793580, y=1600200, w=3886200, h=4572000),
                        text_style={"font_name": "Calibri", "font_size": 14, "font_color": "#333333"},
                        max_chars=300,
                    ),
                ],
                tags=["content", "two_column"],
            ),
            LayoutSpec(
                id="fallback_blank",
                name="Blank",
                layout_type="blank",
                placeholders=[],
                tags=["blank"],
            ),
        ]
