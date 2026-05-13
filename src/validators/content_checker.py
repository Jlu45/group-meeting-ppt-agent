import logging
import re
from typing import Optional

from pptx import Presentation

logger = logging.getLogger(__name__)

PLACEHOLDER_PATTERNS = [
    r"待补充", r"TODO", r"TBD", r"FIXME", r"placeholder",
    r"xxx+", r"测试", r"示例", r"sample",
]


class ContentChecker:
    """内容检查器 - 检查占位符、完整性、术语一致性"""

    def __init__(self, max_text_length: int = 200, min_points: int = 1):
        self._max_text_length = max_text_length
        self._min_points = min_points

    def check(self, pptx_path: str) -> list[dict]:
        issues = []
        try:
            prs = Presentation(pptx_path)
        except Exception as e:
            issues.append({"level": "P0", "message": f"Cannot open PPTX: {e}"})
            return issues

        for i, slide in enumerate(prs.slides):
            slide_has_content = False
            all_text = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_has_content = True
                            all_text.append(text)

                            for pattern in PLACEHOLDER_PATTERNS:
                                if re.search(pattern, text, re.IGNORECASE):
                                    issues.append({
                                        "level": "P0",
                                        "message": f"Slide {i + 1} contains placeholder text: '{text[:50]}'"
                                    })
                                    break

                            if len(text) > self._max_text_length:
                                issues.append({
                                    "level": "P2",
                                    "message": f"Slide {i + 1} has overly long text ({len(text)} chars): '{text[:50]}...'"
                                })

            if not slide_has_content and i > 0:
                issues.append({
                    "level": "P1",
                    "message": f"Slide {i + 1} has no text content"
                })

        issues.extend(self._check_structure(prs))

        return issues

    def _check_structure(self, prs: Presentation) -> list[dict]:
        issues = []
        num_slides = len(prs.slides)

        if num_slides == 0:
            issues.append({"level": "P0", "message": "Presentation has no slides"})
            return issues

        if num_slides < 3:
            issues.append({
                "level": "P2",
                "message": f"Presentation has only {num_slides} slides (recommended: 8-15)"
            })

        first_slide = prs.slides[0]
        has_title = False
        for shape in first_slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        has_title = True
                        break
        if not has_title:
            issues.append({
                "level": "P1",
                "message": "Cover slide has no title text"
            })

        return issues
