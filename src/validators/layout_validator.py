import logging
from typing import Optional

from pptx import Presentation
from pptx.util import Emu

logger = logging.getLogger(__name__)


class LayoutValidator:
    """布局验证器 - 检查溢出、重叠、空白页"""

    def __init__(self, overflow_threshold: float = 0.95, overlap_threshold: float = 0.3):
        self._overflow_threshold = overflow_threshold
        self._overlap_threshold = overlap_threshold

    def validate(self, pptx_path: str) -> list[dict]:
        issues = []
        try:
            prs = Presentation(pptx_path)
        except Exception as e:
            issues.append({"level": "P0", "message": f"Cannot open PPTX: {e}"})
            return issues

        slide_w = prs.slide_width
        slide_h = prs.slide_height

        for i, slide in enumerate(prs.slides):
            shapes_info = []

            for shape in slide.shapes:
                shape_issues = self._check_shape(shape, i, slide_w, slide_h)
                issues.extend(shape_issues)

                shapes_info.append({
                    "name": shape.name,
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                })

            overlap_issues = self._check_overlaps(shapes_info, i)
            issues.extend(overlap_issues)

            if not slide.shapes:
                issues.append({
                    "level": "P1",
                    "message": f"Slide {i + 1} is empty (no shapes)"
                })

        return issues

    def _check_shape(self, shape, slide_idx: int, slide_w: int, slide_h: int) -> list[dict]:
        issues = []

        right = shape.left + shape.width
        bottom = shape.top + shape.height

        if right > slide_w * self._overflow_threshold:
            overflow_pct = (right - slide_w) / slide_w * 100
            issues.append({
                "level": "P1",
                "message": f"Slide {slide_idx + 1} shape '{shape.name}' overflows right by {overflow_pct:.1f}%"
            })

        if bottom > slide_h * self._overflow_threshold:
            overflow_pct = (bottom - slide_h) / slide_h * 100
            issues.append({
                "level": "P1",
                "message": f"Slide {slide_idx + 1} shape '{shape.name}' overflows bottom by {overflow_pct:.1f}%"
            })

        if shape.width < Emu(100000) or shape.height < Emu(100000):
            issues.append({
                "level": "P2",
                "message": f"Slide {slide_idx + 1} shape '{shape.name}' has near-zero dimensions"
            })

        return issues

    def _check_overlaps(self, shapes: list[dict], slide_idx: int) -> list[dict]:
        issues = []
        for i in range(len(shapes)):
            for j in range(i + 1, len(shapes)):
                overlap = self._calculate_overlap(shapes[i], shapes[j])
                if overlap > self._overlap_threshold:
                    issues.append({
                        "level": "P2",
                        "message": f"Slide {slide_idx + 1} shapes '{shapes[i]['name']}' and '{shapes[j]['name']}' overlap by {overlap:.0%}"
                    })
        return issues

    def _calculate_overlap(self, s1: dict, s2: dict) -> float:
        x_overlap = max(0, min(s1["left"] + s1["width"], s2["left"] + s2["width"]) - max(s1["left"], s2["left"]))
        y_overlap = max(0, min(s1["top"] + s1["height"], s2["top"] + s2["height"]) - max(s1["top"], s2["top"]))

        if x_overlap == 0 or y_overlap == 0:
            return 0.0

        intersection = x_overlap * y_overlap
        area1 = s1["width"] * s1["height"]
        area2 = s2["width"] * s2["height"]
        smaller = min(area1, area2)

        if smaller == 0:
            return 0.0

        return intersection / smaller
