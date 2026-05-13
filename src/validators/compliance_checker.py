import logging
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor

from src.models import TemplateDNA, ComplianceReport

logger = logging.getLogger(__name__)


class ComplianceChecker:
    """合规检查器 - 检查PPTX是否符合模板DNA"""

    def __init__(self, template_dna: Optional[TemplateDNA] = None):
        self._template_dna = template_dna or TemplateDNA()

    def check(self, pptx_path: str) -> ComplianceReport:
        try:
            prs = Presentation(pptx_path)
        except Exception as e:
            logger.error(f"Cannot open PPTX for compliance check: {e}")
            return ComplianceReport()

        color_score = self._check_color_usage(prs)
        font_score = self._check_font_usage(prs)
        layout_score = self._check_layout_consistency(prs)
        decoration_score = self._check_decorations(prs)

        return ComplianceReport(
            color_score=color_score,
            font_score=font_score,
            layout_score=layout_score,
            decoration_score=decoration_score,
        )

    def _check_color_usage(self, prs: Presentation) -> float:
        theme = self._template_dna.theme
        allowed_colors = self._get_allowed_colors(theme)

        total_runs = 0
        compliant_runs = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            try:
                                if run.font.color and run.font.color.type is not None:
                                    color_str = str(run.font.color.rgb).upper().lstrip("#")
                                    total_runs += 1
                                    if color_str in allowed_colors:
                                        compliant_runs += 1
                            except (AttributeError, TypeError):
                                pass

        if total_runs == 0:
            return 1.0

        return compliant_runs / total_runs

    def _check_font_usage(self, prs: Presentation) -> float:
        fonts = self._template_dna.fonts
        allowed_fonts = {fonts.title, fonts.body, fonts.mono}

        total_runs = 0
        compliant_runs = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name:
                                total_runs += 1
                                if run.font.name in allowed_fonts:
                                    compliant_runs += 1

        if total_runs == 0:
            return 1.0

        return compliant_runs / total_runs

    def _check_layout_consistency(self, prs: Presentation) -> float:
        if len(prs.slides) < 2:
            return 1.0

        title_positions = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size and run.font.size >= 200000:
                                title_positions.append(shape.top)
                                break

        if len(title_positions) < 2:
            return 0.9

        avg_pos = sum(title_positions) / len(title_positions)
        deviations = [abs(p - avg_pos) for p in title_positions]
        max_deviation = max(deviations) if deviations else 0

        if max_deviation < 200000:
            return 1.0
        elif max_deviation < 500000:
            return 0.9
        else:
            return 0.7

    def _check_decorations(self, prs: Presentation) -> float:
        decorations = self._template_dna.decorations
        if not decorations or (not decorations.logo and not decorations.footer):
            return 0.9

        return 0.9

    def _get_allowed_colors(self, theme) -> set:
        colors = set()
        for attr in ["primary", "secondary", "accent", "background", "text", "light_bg", "subtle"]:
            val = getattr(theme, attr, "")
            if val:
                colors.add(val.upper().lstrip("#"))

        colors.update({"FFFFFF", "000000", "DDDDDD", "BBBBBB"})
        return colors
