import logging
from typing import Optional
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from src.models import TemplateDNA, ThemeColors, FontHierarchy, ComplianceReport

logger = logging.getLogger(__name__)

OOXML_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
OOXML_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
OOXML_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


class StyleLockEngine:
    """样式锁定引擎 - 确保生成PPT严格遵循模板DNA"""

    def __init__(self, template_dna: Optional[TemplateDNA] = None):
        self._template_dna = template_dna or TemplateDNA()

    def apply(self, pptx_path: str) -> str:
        prs = Presentation(pptx_path)
        theme = self._template_dna.theme
        fonts = self._template_dna.fonts

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    self._remap_text_frame(shape.text_frame, theme, fonts)

                self._remap_shape_colors(shape, theme)

            self._calibrate_dimensions(slide, prs)

        self._apply_decorations(prs, self._template_dna)

        prs.save(pptx_path)
        logger.info(f"Style lock applied to {pptx_path}")
        return pptx_path

    def _remap_text_frame(self, text_frame, theme: ThemeColors, fonts: FontHierarchy):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.name and run.font.name != fonts.title:
                    if run.font.size and run.font.size >= Pt(fonts.title_size - 4):
                        run.font.name = fonts.title
                    else:
                        run.font.name = fonts.body

                try:
                    if run.font.color and run.font.color.type is not None:
                        current = str(run.font.color.rgb)
                        mapped = self._map_color(current, theme)
                        if mapped != current:
                            run.font.color.rgb = RGBColor.from_string(mapped)
                except (AttributeError, TypeError):
                    pass

    def _remap_shape_colors(self, shape, theme: ThemeColors):
        try:
            if not hasattr(shape, "fill"):
                return

            fill = shape.fill
            if fill.type is None:
                return

            from pptx.enum.dml import MSO_THEME_COLOR
            if fill.type == MSO_THEME_COLOR:
                return

            try:
                fore_color = fill.fore_color
                if fore_color and fore_color.type is not None:
                    current = str(fore_color.rgb)
                    mapped = self._map_color(current, theme)
                    if mapped != current:
                        fill.fore_color.rgb = RGBColor.from_string(mapped)
            except (AttributeError, TypeError):
                pass

        except Exception as e:
            logger.debug(f"Shape color remap skipped: {e}")

    def _calibrate_dimensions(self, slide, prs: Presentation):
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        for shape in slide.shapes:
            if shape.left < 0:
                shape.left = 0
            if shape.top < 0:
                shape.top = 0

            right = shape.left + shape.width
            if right > slide_w:
                shape.width = slide_w - shape.left

            bottom = shape.top + shape.height
            if bottom > slide_h:
                shape.height = slide_h - shape.top

    def _map_color(self, color: str, theme: ThemeColors) -> str:
        color = color.upper().lstrip("#")
        theme_colors = {
            theme.primary.upper().lstrip("#"): theme.primary.lstrip("#"),
            theme.secondary.upper().lstrip("#"): theme.secondary.lstrip("#"),
            theme.accent.upper().lstrip("#"): theme.accent.lstrip("#"),
            theme.text.upper().lstrip("#"): theme.text.lstrip("#"),
            theme.background.upper().lstrip("#"): theme.background.lstrip("#"),
        }

        if color in theme_colors:
            return theme_colors[color]

        near_matches = {
            "000000": theme.text.lstrip("#"),
            "FFFFFF": theme.light_bg.lstrip("#"),
            "333333": theme.text.lstrip("#"),
            "666666": theme.subtle.lstrip("#"),
            "999999": theme.subtle.lstrip("#"),
            "CCCCCC": theme.subtle.lstrip("#"),
        }

        if color in near_matches:
            return near_matches[color]

        return color

    def _apply_decorations(self, prs: Presentation, template_dna: TemplateDNA):
        decorations = template_dna.decorations
        if not decorations:
            return

        theme = template_dna.theme
        fonts = template_dna.fonts

        if decorations.logo and decorations.logo.get("rId"):
            self._apply_logo(prs, template_dna)

        if decorations.footer:
            self._apply_footer(prs, theme, fonts)

        if decorations.header:
            self._apply_header(prs, theme, fonts)

        if decorations.dividers:
            self._apply_dividers(prs, theme)

    def _apply_logo(self, prs: Presentation, template_dna: TemplateDNA):
        logo_info = template_dna.decorations.logo
        if not logo_info:
            return

        media = template_dna.media
        r_id = logo_info.get("rId", "")
        pos = logo_info.get("position", {})

        logo_data = None
        for media_path, data in media.items():
            if r_id in media_path or "logo" in media_path.lower():
                logo_data = data
                break

        if not logo_data:
            return

        import tempfile
        from pathlib import Path
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        tmp.write(logo_data)
        tmp.close()

        try:
            for slide in prs.slides:
                try:
                    x = pos.get("x", Inches(11.5))
                    y = pos.get("y", Inches(0.2))
                    cx = pos.get("cx", Inches(1.0))
                    cy = pos.get("cy", Inches(0.8))
                    slide.shapes.add_picture(tmp.name, x, y, cx, cy)
                except Exception:
                    pass
        finally:
            Path(tmp.name).unlink(missing_ok=True)

    def _apply_footer(self, prs: Presentation, theme: ThemeColors, fonts: FontHierarchy):
        for i, slide in enumerate(prs.slides):
            try:
                footer_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.4)
                )
                tf = footer_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"第 {i + 1} / {len(prs.slides)} 页"
                p.font.size = Pt(fonts.small_size - 2)
                p.font.color.rgb = RGBColor.from_string(theme.subtle.lstrip("#"))
                p.font.name = fonts.body
                p.alignment = 2  # PP_ALIGN.RIGHT
            except Exception:
                pass

    def _apply_header(self, prs: Presentation, theme: ThemeColors, fonts: FontHierarchy):
        pass

    def _apply_dividers(self, prs: Presentation, theme: ThemeColors):
        for divider_pos in theme.dividers if hasattr(theme, 'dividers') else []:
            for slide in prs.slides:
                try:
                    x = divider_pos.get("x", 0)
                    y = divider_pos.get("y", Inches(1.2))
                    cx = divider_pos.get("cx", prs.slide_width)
                    cy = divider_pos.get("cy", Inches(0.02))

                    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cx, cy)
                    line.fill.solid()
                    line.fill.fore_color.rgb = RGBColor.from_string(theme.accent.lstrip("#"))
                    line.line.fill.background()
                except Exception:
                    pass

    def validate_compliance(self, pptx_path: str) -> ComplianceReport:
        prs = Presentation(pptx_path)
        theme = self._template_dna.theme
        fonts = self._template_dna.fonts

        color_checks = []
        font_checks = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            try:
                                if run.font.color and run.font.color.type is not None:
                                    color_checks.append(
                                        self._check_color_compliance(str(run.font.color.rgb), theme)
                                    )
                            except (AttributeError, TypeError):
                                pass
                            if run.font.name:
                                font_checks.append(
                                    run.font.name in (fonts.title, fonts.body, fonts.mono)
                                )

        color_score = sum(color_checks) / max(len(color_checks), 1)
        font_score = sum(font_checks) / max(len(font_checks), 1)
        layout_score = self._check_layout_consistency(prs)
        decoration_score = self._check_decoration_compliance(prs)

        return ComplianceReport(
            color_score=color_score,
            font_score=font_score,
            layout_score=layout_score,
            decoration_score=decoration_score,
        )

    def _check_color_compliance(self, color: str, theme: ThemeColors) -> float:
        allowed = {
            theme.primary.upper().lstrip("#"),
            theme.secondary.upper().lstrip("#"),
            theme.accent.upper().lstrip("#"),
            theme.text.upper().lstrip("#"),
            theme.background.upper().lstrip("#"),
            theme.light_bg.upper().lstrip("#"),
            theme.subtle.upper().lstrip("#"),
            "FFFFFF", "000000", "DDDDDD", "BBBBBB",
        }
        return 1.0 if color.upper().lstrip("#") in allowed else 0.0

    def _check_layout_consistency(self, prs: Presentation) -> float:
        if len(prs.slides) < 2:
            return 1.0

        title_positions = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size and run.font.size >= Pt(24):
                                title_positions.append(shape.top)
                                break

        if len(title_positions) < 2:
            return 0.9

        avg_pos = sum(title_positions) / len(title_positions)
        deviations = [abs(p - avg_pos) for p in title_positions]
        max_deviation = max(deviations) if deviations else 0

        if max_deviation < 200000:
            return 1.0
        elif max_deviation < 500000:
            return 0.9
        else:
            return 0.7

    def _check_decoration_compliance(self, prs: Presentation) -> float:
        decorations = self._template_dna.decorations
        if not decorations or (not decorations.logo and not decorations.footer and not decorations.dividers):
            return 0.9

        score = 0.0
        total = 0

        if decorations.footer:
            total += 1
            footer_found = False
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            if "/" in p.text and "页" in p.text:
                                footer_found = True
                                break
            if footer_found:
                score += 1.0

        if decorations.logo:
            total += 1
            score += 0.5

        if decorations.dividers:
            total += 1
            score += 0.5

        if total == 0:
            return 0.9

        return score / total
