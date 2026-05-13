import io
import logging
import tempfile
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from src.models import (
    StructuredPresentation, SlidePlan, LayoutType, TemplateDNA,
    ThemeColors, FontHierarchy,
)

logger = logging.getLogger(__name__)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


class PPTXBuilder:
    """PPTX构建器 - 将结构化数据转化为原生可编辑的PPTX"""

    def __init__(self, template_dna: Optional[TemplateDNA] = None, template_path: Optional[str] = None, chart_generator=None):
        self._template_dna = template_dna or TemplateDNA()
        self._template_path = template_path
        self._chart_generator = chart_generator
        self._temp_files = []

    def build(self, presentation: StructuredPresentation, output_path: str) -> str:
        prs = self._create_presentation()

        for i, slide_plan in enumerate(presentation.slides):
            self._build_slide(prs, slide_plan, presentation, i)

        output_dir = Path(output_path).parent
        output_dir.mkdir(parents=True, exist_ok=True)

        prs.save(output_path)
        logger.info(f"PPTX saved to {output_path}")

        for tf in self._temp_files:
            try:
                Path(tf).unlink(missing_ok=True)
            except Exception:
                pass
        self._temp_files.clear()

        return output_path

    def _create_presentation(self) -> Presentation:
        if self._template_path and Path(self._template_path).exists():
            prs = Presentation(self._template_path)
            logger.info("Using template PPTX as base")
            return prs

        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        return prs

    def _build_slide(self, prs: Presentation, plan: SlidePlan, presentation: StructuredPresentation, index: int):
        layout_map = {
            LayoutType.COVER: self._build_cover_slide,
            LayoutType.BULLET_LIST: self._build_bullet_slide,
            LayoutType.TWO_COLUMN: self._build_two_column_slide,
            LayoutType.CHART: self._build_chart_slide,
            LayoutType.TABLE: self._build_table_slide,
            LayoutType.IMAGE_GRID: self._build_image_grid_slide,
            LayoutType.ARCHITECTURE: self._build_architecture_slide,
            LayoutType.SUMMARY: self._build_summary_slide,
            LayoutType.DISCUSSION: self._build_discussion_slide,
        }

        builder = layout_map.get(plan.layout, self._build_bullet_slide)
        slide = builder(prs, plan, presentation, index)

        if plan.notes:
            self._add_notes(slide, plan.notes)

    def _build_cover_slide(self, prs: Presentation, plan: SlidePlan, presentation: StructuredPresentation, index: int):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = self._template_dna.theme
        fonts = self._template_dna.fonts

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(theme.primary.lstrip("#"))

        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = plan.title or presentation.title
        p.font.size = Pt(fonts.title_size + 8)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        p.font.name = fonts.title
        p.alignment = PP_ALIGN.CENTER

        if presentation.author:
            author_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.3), Inches(11.333), Inches(0.6)
            )
            tf = author_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"汇报人：{presentation.author}"
            p.font.size = Pt(fonts.body_size)
            p.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            p.font.name = fonts.body
            p.alignment = PP_ALIGN.CENTER

        if presentation.date:
            date_box = slide.shapes.add_textbox(
                Inches(1), Inches(5.0), Inches(11.333), Inches(0.5)
            )
            tf = date_box.text_frame
            p = tf.paragraphs[0]
            p.text = presentation.date
            p.font.size = Pt(fonts.small_size)
            p.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
            p.font.name = fonts.body
            p.alignment = PP_ALIGN.CENTER

        self._add_accent_line(slide, Inches(4.0), Inches(4.1), Inches(5.333))

        return slide

    def _build_bullet_slide(self, prs: Presentation, plan: SlidePlan, presentation: StructuredPresentation, index: int):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = self._template_dna.theme
        fonts = self._template_dna.fonts

        self._set_slide_bg(slide, theme.background)
        self._add_title_bar(slide, plan.title, theme, fonts)

        content_top = Inches(1.6)
        content_left = Inches(1.0)
        content_width = Inches(11.333)
        content_height = Inches(5.4)

        body_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
        tf = body_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(plan.points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(fonts.body_size)
            p.font.color.rgb = RGBColor.from_string(theme.text.lstrip("#"))
            p.font.name = fonts.body
            p.space_after = Pt(10)
            p.level = 0

        return slide

    def _build_two_column_slide(self, prs: Presentation, plan: SlidePlan, presentation: StructuredPresentation, index: int):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = self._template_dna.theme
        fonts = self._template_dna.fonts

        self._set_slide_bg(slide, theme.background)
        self._add_title_bar(slide, plan.title, theme, fonts)

        mid = len(plan.points) // 2
        left_points = plan.points[:mid] if mid > 0 else plan.points
        right_points = plan.points[mid:] if mid > 0 else []

        col_width = Inches(5.4)
        left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), col_width, Inches(5.4))
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(left_points):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(fonts.body_size)
            p.font.color.rgb = RGBColor.from_string(theme.text.lstrip("#"))
            p.font.name = fonts.body
            p.space_after = Pt(8)

        if right_points:
            right_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), col_width, Inches(5.4))
            tf = right_box.text_frame
            tf.word_wrap = True
            for i, point in enumerate(right_points):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(fonts.body_size)
                p.font.color.rgb = RGBColor.from_string(theme.text.lstrip("#"))
                p.font.name = fonts.body
                p.space_after = Pt(8)

        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.8), Inches(0.03), Inches(5.0)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = RGBColor.from_string(theme.subtle.lstrip("#"))
        divider.line.fill.background()

        return slide

    def _build_chart_slide(self, prs: Presentation, plan: SlidePlan, presentation: StructuredPresentation, index: int):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = self._template_dna.theme
        fonts = self._template_dna.fonts

        self._set_slide_bg(slide, theme.background)
        self._add_title_bar(slide, plan.title, theme, fonts)

        chart_image_path = self._generate_chart_image(plan, index)

        if chart_image_path:
            slide.shapes.add_picture(
                chart_image_path,
                Inches(1.5), Inches(1.6),
                Inches(10.0), Inches(4.5),
            )
        else:
            placeholder_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.8), Inches(11.333), Inches(4.5)
            )
            tf = placeholder_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"[图表区域] {plan.chart_desc or '待生成图表'}"
            p.font.size = Pt(fonts.body_size)
            p.font.color.rgb = RGBColor.from_string(theme.subtle.lstrip("#"))
            p.font.name = fonts.body
            p.alignment = PP_ALIGN.CENTER

        if plan.points:
            notes_box = slide.shapes.add_textbox(
                Inches(1), Inches(6.3), Inches(11.333), Inches(0.8)
            )
            tf = notes_box.text_frame
            tf.word_wrap = True
            for i, point in enumerate(plan.points[:3]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(fonts.small_size)
                p.font.color.rgb = RGBColor.from_string(theme.text.lstrip("#"))
                p.font.name = fonts.body

        return slide

    def _generate_chart_image(self, plan: SlidePlan, index: int) -> Optional[str]:
        if not self._chart_generator or not plan.chart_desc:
            return None

        try:
            tmpdir = tempfile.mkdtemp()
            chart_path = str(Path(tmpdir) / f"chart_{index}.png")

            if plan.table_data and len(plan.table_data) > 1:
                headers = plan.table_data[0]
                data = plan.table_data[1:]
                self._chart_generator.generate_comparison_table(
                    data=data, headers=headers, output_path=chart_path,
                )
            else:
                self._chart_generator.generate_bar_chart(
                    labels=[f"Item {i + 1}" for i in range(len(plan.points))],
                    values=list(range(1, len(plan.points) + 1)),
                    title=plan.chart_desc,
                    output_path=chart_path,
                )

            if Path(chart_path).exists() and Path(chart_path).stat().st_size > 0:
                self._temp_files.append(chart_path)
                return chart_path
        except Exception as e:
            logger.warning(f"Chart generation failed for slide '{plan.title}': {e}")

        return None

    def _build_table_slide(self, prs: Presentation, plan: SlidePlan, presentation: StructuredPresentation, index: int):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = self._template_dna.theme
        fonts = self._template_dna.fonts

        self._set_slide_bg(slide, theme.background)
        self._add_title_bar(slide, plan.title, theme, fonts)

        if plan.table_data and len(plan.table_data) > 0:
            headers = plan.table_data[0] if plan.table_data else []
            rows = plan.table_data[1:] if len(plan.table_data) > 1 else []

            num_rows = len(rows) + 1
            num_cols = len(headers) if headers else 1

            table_width = Inches(11.333)
            table_height = Inches(4.5)
            table_shape = slide.shapes.add_table(
                num_rows, num_cols,
                Inches(1), Inches(1.8),
                table_width, table_height,
            )
            table = table_shape.table

            for j, header in enumerate(headers):
                cell = table.cell(0, j)
                cell.text = str(header)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(fonts.small_size)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    paragraph.font.name = fonts.body
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(theme.primary.lstrip("#"))

            for i, row in enumerate(rows):
                for j, val in enumerate(row):
                    if j < num_cols:
                        cell = table.cell(i + 1, j)
                        cell.text = str(val)
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(fonts.small_size)
                            paragraph.font.color.rgb = RGBColor.from_string(theme.text.lstrip("#"))
                            paragraph.font.name = fonts.body
                        if i % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor.from_string(theme.light_bg.lstrip("#"))
        else:
            body_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.8), Inches(11.333), Inches(4.5)
            )
            tf = body_box.text_frame
            tf.word_wrap = True
            for i, point in enumerate(plan.points):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(fonts.body_size)
                p.font.color.rgb = RGBColor.from_string(theme.text.lstrip("#"))
                p.font.name = fonts.body
                p.space_after = Pt(8)

        return slide

    def _build_image_grid_slide(self, prs: Presentation, plan: SlidePlan, presentation: StructuredPresentation, index: int):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = self._template_dna.theme
        fonts = self._template_dna.fonts

        self._set_slide_bg(slide, theme.background)
        self._add_title_bar(slide, plan.title, theme, fonts)

        valid_images = [p for p in plan.image_paths if Path(p).exists()]

        if valid_images:
            num_images = len(valid_images)
            if num_images <= 2:
                cols, rows = 2, 1
            elif num_images <= 4:
                cols, rows = 2, 2
            elif num_images <= 6:
                cols, rows = 3, 2
            else:
                cols, rows = 4, 2

            content_top = Inches(1.5)
            content_left = Inches(0.8)
            total_w = Inches(11.733)
            total_h = Inches(5.5)
            gap = Inches(0.2)
            cell_w = (total_w - gap * (cols - 1)) / cols
            cell_h = (total_h - gap * (rows - 1)) / rows

            for idx, img_path in enumerate(valid_images[:cols * rows]):
                r = idx // cols
                c = idx % cols
                x = content_left + c * (cell_w + gap)
                y = content_top + r * (cell_h + gap)
                try:
                    slide.shapes.add_picture(img_path, x, y, cell_w, cell_h)
                except Exception as e:
                    logger.warning(f"Failed to insert image {img_path}: {e}")
        else:
            body_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.8), Inches(11.333), Inches(5.0)
            )
            tf = body_box.text_frame
            tf.word_wrap = True
            for i, point in enumerate(plan.points):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(fonts.body_size)
                p.font.color.rgb = RGBColor.from_string(theme.text.lstrip("#"))
                p.font.name = fonts.body
                p.space_after = Pt(8)

        return slide

    def _build_architecture_slide(self, prs: Presentation, plan: SlidePlan, presentation: StructuredPresentation, index: int):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = self._template_dna.theme
        fonts = self._template_dna.fonts

        self._set_slide_bg(slide, theme.background)
        self._add_title_bar(slide, plan.title, theme, fonts)

        arch_image_path = self._generate_architecture_image(plan, index)

        if arch_image_path:
            slide.shapes.add_picture(
                arch_image_path,
                Inches(1.5), Inches(1.6),
                Inches(10.0), Inches(5.0),
            )
        else:
            self._draw_architecture_blocks(slide, plan, theme, fonts)

        return slide

    def _generate_architecture_image(self, plan: SlidePlan, index: int) -> Optional[str]:
        if not self._chart_generator:
            return None

        try:
            tmpdir = tempfile.mkdtemp()
            arch_path = str(Path(tmpdir) / f"arch_{index}.png")
            self._chart_generator.generate_architecture_diagram(
                components=plan.points,
                title=plan.title,
                output_path=arch_path,
            )
            if Path(arch_path).exists() and Path(arch_path).stat().st_size > 0:
                self._temp_files.append(arch_path)
                return arch_path
        except Exception as e:
            logger.warning(f"Architecture diagram generation failed: {e}")

        return None

    def _draw_architecture_blocks(self, slide, plan: SlidePlan, theme: ThemeColors, fonts: FontHierarchy):
        if not plan.points:
            return

        n = len(plan.points)
        block_w = Inches(2.2)
        block_h = Inches(1.0)
        start_x = Inches(0.8)
        start_y = Inches(2.0)
        gap_x = Inches(0.5)
        gap_y = Inches(0.6)
        cols = min(n, 4)

        for i, point in enumerate(plan.points):
            row = i // cols
            col = i % cols
            x = start_x + col * (block_w + gap_x)
            y = start_y + row * (block_h + gap_y)

            block = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, y, block_w, block_h,
            )
            block.fill.solid()
            if row == 0:
                block.fill.fore_color.rgb = RGBColor.from_string(theme.primary.lstrip("#"))
            else:
                block.fill.fore_color.rgb = RGBColor.from_string(theme.secondary.lstrip("#"))
            block.line.fill.background()

            tf = block.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            p = tf.paragraphs[0]
            p.text = point
            p.font.size = Pt(fonts.small_size)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.bold = True
            p.font.name = fonts.body

            if col > 0 and i > 0:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    x - gap_x + Inches(0.05), y + block_h // 2 - Inches(0.15),
                    gap_x - Inches(0.1), Inches(0.3),
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = RGBColor.from_string(theme.accent.lstrip("#"))
                arrow.line.fill.background()

    def _build_summary_slide(self, prs: Presentation, plan: SlidePlan, presentation: StructuredPresentation, index: int):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = self._template_dna.theme
        fonts = self._template_dna.fonts

        self._set_slide_bg(slide, theme.background)

        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = plan.title or "总结"
        p.font.size = Pt(fonts.title_size)
        p.font.color.rgb = RGBColor.from_string(theme.primary.lstrip("#"))
        p.font.bold = True
        p.font.name = fonts.title
        p.alignment = PP_ALIGN.CENTER

        self._add_accent_line(slide, Inches(5.5), Inches(1.4), Inches(2.333))

        body_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(1.8), Inches(10.333), Inches(5.0)
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(plan.points):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"✓ {point}"
            p.font.size = Pt(fonts.body_size + 2)
            p.font.color.rgb = RGBColor.from_string(theme.text.lstrip("#"))
            p.font.name = fonts.body
            p.space_after = Pt(14)

        return slide

    def _build_discussion_slide(self, prs: Presentation, plan: SlidePlan, presentation: StructuredPresentation, index: int):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = self._template_dna.theme
        fonts = self._template_dna.fonts

        self._set_slide_bg(slide, theme.background)

        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = plan.title or "讨论与下一步"
        p.font.size = Pt(fonts.title_size)
        p.font.color.rgb = RGBColor.from_string(theme.primary.lstrip("#"))
        p.font.bold = True
        p.font.name = fonts.title
        p.alignment = PP_ALIGN.CENTER

        self._add_accent_line(slide, Inches(5.5), Inches(1.4), Inches(2.333))

        body_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(1.8), Inches(10.333), Inches(5.0)
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(plan.points):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"→ {point}"
            p.font.size = Pt(fonts.body_size + 2)
            p.font.color.rgb = RGBColor.from_string(theme.text.lstrip("#"))
            p.font.name = fonts.body
            p.space_after = Pt(14)

        return slide

    def _set_slide_bg(self, slide, color_str: str):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(color_str.lstrip("#"))

    def _add_title_bar(self, slide, title: str, theme: ThemeColors, fonts: FontHierarchy):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.2)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(theme.primary.lstrip("#"))
        bar.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.15), Inches(11.533), Inches(0.9)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(fonts.subtitle_size)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        p.font.name = fonts.title
        p.alignment = PP_ALIGN.LEFT

    def _add_accent_line(self, slide, left, top, width):
        theme = self._template_dna.theme
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor.from_string(theme.accent.lstrip("#"))
        line.line.fill.background()

    def _add_notes(self, slide, notes: str):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
