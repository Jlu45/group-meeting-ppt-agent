import os
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))

from src.models import (
    DocType, LayoutType, ParsedDocument, TemplateDNA, ThemeColors,
    FontHierarchy, SlidePlan, StructuredPresentation, ValidationResult,
    ComplianceReport,
)
from src.parsers.document_parser import UniversalDocumentParser
from src.parsers.template_extractor import TemplateDNAExtractor
from src.parsers.content_structurer import ContentStructuringEngine
from src.agents.planner import PlannerAgent
from src.agents.generator import GeneratorAgent
from src.agents.refiner import RefinerAgent
from src.agents.validator import ValidatorAgent
from src.generators.pptx_builder import PPTXBuilder
from src.generators.chart_generator import ChartGenerator
from src.generators.style_lock import StyleLockEngine
from src.generators.layout_engine import LayoutEngine
from src.validators.layout_validator import LayoutValidator
from src.validators.compliance_checker import ComplianceChecker
from src.validators.content_checker import ContentChecker


def test_models():
    print("Testing models...")
    dna = TemplateDNA()
    assert dna.theme.primary == "#1E3A5F"
    assert dna.fonts.title == "Source Han Serif SC"

    report = ComplianceReport(color_score=0.9, font_score=0.95, layout_score=0.85, decoration_score=0.9)
    assert report.overall_score == 0.9
    assert report.passed

    result = ValidationResult(level1_passed=True, level2_passed=True, level3_passed=True)
    assert result.all_passed
    print("  ✓ Models test passed")


def test_document_parser():
    print("Testing document parser...")
    parser = UniversalDocumentParser()

    with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False, encoding="utf-8") as f:
        f.write("# Test Document\n\nThis is a test.\n\n- Point 1\n- Point 2\n")
        temp_path = f.name

    try:
        doc = parser.parse(temp_path)
        assert doc.source_path == str(Path(temp_path).resolve())
        assert "Test Document" in doc.markdown_content
        assert doc.file_type == "md"
        print("  ✓ Document parser test passed")
    finally:
        os.unlink(temp_path)


def test_content_structurer():
    print("Testing content structurer...")
    structurer = ContentStructuringEngine(llm_client=None)

    doc = ParsedDocument(
        source_path="test.md",
        markdown_content="# 实验结果\n\n## 方法\n我们使用了XX方法\n\n## 结果\n准确率达到95%\n\n## 分析\n结果优于基线",
        file_type="md",
    )

    presentation = structurer.structure(doc, author="张三", date="2025-01-01")
    assert presentation.title == "实验结果"
    assert len(presentation.slides) >= 4
    assert presentation.slides[0].layout == LayoutType.COVER
    assert any(s.layout == LayoutType.SUMMARY for s in presentation.slides)
    assert any(s.layout == LayoutType.DISCUSSION for s in presentation.slides)
    print("  ✓ Content structurer test passed")


def test_planner():
    print("Testing planner...")
    planner = PlannerAgent(llm_client=None)

    presentation = StructuredPresentation(
        doc_type=DocType.EXPERIMENT_LOG,
        title="实验结果",
        summary="测试",
        slides=[
            SlidePlan(title="封面", layout=LayoutType.COVER, points=[]),
            SlidePlan(title="概述", layout=LayoutType.BULLET_LIST, points=["点1", "点2", "点3"]),
            SlidePlan(title="方法", layout=LayoutType.BULLET_LIST, points=["a", "b", "c", "d", "e", "f", "g"]),
            SlidePlan(title="总结", layout=LayoutType.SUMMARY, points=["结论1"]),
            SlidePlan(title="讨论", layout=LayoutType.DISCUSSION, points=["下一步"]),
        ],
    )

    result = planner.plan(presentation)
    assert len(result.slides) >= 4
    assert result.slides[0].layout == LayoutType.COVER
    print("  ✓ Planner test passed")


def test_refiner():
    print("Testing refiner...")
    refiner = RefinerAgent(llm_client=None)

    presentation = StructuredPresentation(
        doc_type=DocType.OTHER,
        title="测试",
        summary="",
        slides=[
            SlidePlan(title="封面", layout=LayoutType.COVER, points=[]),
            SlidePlan(title="内容", layout=LayoutType.BULLET_LIST,
                      points=["短", "这是一个非常非常非常非常非常非常非常非常非常非常非常非常非常非常长的要点需要截断处理"]),
            SlidePlan(title="总结", layout=LayoutType.SUMMARY, points=["结论"]),
            SlidePlan(title="讨论", layout=LayoutType.DISCUSSION, points=["下一步"]),
        ],
    )

    result = refiner.refine(presentation)
    long_point = result.slides[1].points[1]
    assert len(long_point) <= 63
    print("  ✓ Refiner test passed")


def test_validator():
    print("Testing validator...")
    validator = ValidatorAgent()

    presentation = StructuredPresentation(
        doc_type=DocType.OTHER,
        title="测试",
        summary="",
        slides=[
            SlidePlan(title="封面", layout=LayoutType.COVER, points=[]),
            SlidePlan(title="内容", layout=LayoutType.BULLET_LIST, points=["要点1", "要点2"]),
            SlidePlan(title="总结", layout=LayoutType.SUMMARY, points=["结论"]),
            SlidePlan(title="讨论", layout=LayoutType.DISCUSSION, points=["下一步"]),
        ],
    )

    result = validator.validate_structure(presentation)
    assert result.level1_passed
    print("  ✓ Validator test passed")


def test_pptx_builder():
    print("Testing PPTX builder...")
    builder = PPTXBuilder()

    presentation = StructuredPresentation(
        doc_type=DocType.EXPERIMENT_LOG,
        title="实验结果汇报",
        summary="测试实验结果",
        author="张三",
        date="2025-01-01",
        slides=[
            SlidePlan(title="实验结果汇报", layout=LayoutType.COVER, points=[]),
            SlidePlan(title="概述", layout=LayoutType.BULLET_LIST, points=["实验目标", "核心发现"]),
            SlidePlan(title="方法对比", layout=LayoutType.TWO_COLUMN, points=["方法A: 准确率90%", "方法B: 准确率85%"]),
            SlidePlan(title="结果", layout=LayoutType.TABLE,
                      points=["详见表格"],
                      table_data=[["方法", "准确率", "F1"], ["Ours", "95%", "0.93"], ["Baseline", "88%", "0.85"]]),
            SlidePlan(title="总结", layout=LayoutType.SUMMARY, points=["方法有效", "性能提升7%"]),
            SlidePlan(title="讨论与下一步", layout=LayoutType.DISCUSSION, points=["扩展数据集", "优化超参数"]),
        ],
    )

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = os.path.join(tmpdir, "test_output.pptx")
        builder.build(presentation, output_path)
        assert os.path.exists(output_path)
        assert os.path.getsize(output_path) > 0

        from pptx import Presentation
        prs = Presentation(output_path)
        assert len(prs.slides) == 6

    print("  ✓ PPTX builder test passed")


def test_chart_generator():
    print("Testing chart generator...")
    generator = ChartGenerator()

    with tempfile.TemporaryDirectory() as tmpdir:
        bar_path = os.path.join(tmpdir, "bar.png")
        generator.generate_bar_chart(
            labels=["A", "B", "C"], values=[10, 25, 15],
            title="Test Chart", output_path=bar_path,
        )
        assert os.path.exists(bar_path)
        assert os.path.getsize(bar_path) > 0

        sota_path = os.path.join(tmpdir, "sota.png")
        generator.generate_sota_comparison(
            data=[["Method A", 85.2], ["Method B", 90.1], ["Ours", 95.3]],
            output_path=sota_path,
        )
        assert os.path.exists(sota_path)

    print("  ✓ Chart generator test passed")


def test_layout_engine():
    print("Testing layout engine...")
    engine = LayoutEngine()

    cover_plan = SlidePlan(title="封面", layout=LayoutType.COVER, points=[])
    layout = engine.calculate_layout(cover_plan)
    assert layout["type"] == "cover"
    assert "title" in layout

    bullet_plan = SlidePlan(title="内容", layout=LayoutType.BULLET_LIST, points=["a", "b"])
    layout = engine.calculate_layout(bullet_plan)
    assert layout["type"] == "bullet_list"
    assert "content" in layout

    print("  ✓ Layout engine test passed")


def test_end_to_end():
    print("Testing end-to-end flow...")
    from src.agent import GroupMeetingPPTAgent, GenerationConfig

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "experiment.md")
        with open(input_path, "w", encoding="utf-8") as f:
            f.write("""# XX实验结果汇报

## 实验目标
验证新方法在XX数据集上的有效性

## 实验方法
- 使用Transformer架构
- 数据集：CIFAR-100
- 训练200个epoch

## 主实验结果
| 方法 | 准确率 | F1 |
|------|--------|-----|
| Baseline | 85.2% | 0.83 |
| Ours | 92.1% | 0.91 |

## 消融实验
- 去掉模块A：准确率下降3%
- 去掉模块B：准确率下降5%
- 完整模型：92.1%

## 总结
- 新方法在XX数据集上达到SOTA
- 关键模块A和B均有效

## 下一步
- 在更多数据集上验证
- 优化推理速度
""")

        config = GenerationConfig(
            author="张三",
            date="2025-01-01",
            output_dir=os.path.join(tmpdir, "output"),
            skip_llm=True,
        )

        agent = GroupMeetingPPTAgent(llm_client=None, config=config)
        result = agent.generate([input_path], config)

        assert os.path.exists(result.pptx_path)
        assert len(result.presentation.slides) >= 4
        assert result.duration_seconds > 0

        from pptx import Presentation
        prs = Presentation(result.pptx_path)
        assert len(prs.slides) >= 4

    print("  ✓ End-to-end test passed")


if __name__ == "__main__":
    test_models()
    test_document_parser()
    test_content_structurer()
    test_planner()
    test_refiner()
    test_validator()
    test_pptx_builder()
    test_chart_generator()
    test_layout_engine()
    test_end_to_end()
    print("\n✅ All tests passed!")
